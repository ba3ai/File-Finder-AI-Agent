# smart_content.py
# -----------------------------------------------------------------------------
# Content-first resolver for your file agent (Ollama via langchain_ollama).
#
# NEW: Local upload support
#  - If a hit points to a locally uploaded file (source == "upload"/"local"/"file"
#    or has a path/url under /uploads), this module now reads it directly from disk.
#
# What this module does (dynamic, no domain hardcoding):
#   1) Builds a candidate pool from search hits (or all hits, toggled by env)
#   2) Downloads/reads candidates in parallel (bounded), extracts text (with OCR)
#   3) Ranks with BM25 (length-normalized, IDF-weighted) + rare-term gate
#   4) Picks the best document, narrows to relevant passages
#   5) Calls local Ollama (ChatOllama) to answer STRICTLY from passages
#
# Optional Exhaustive Mode:
#   - Scan ALL files (bounded by caps), chunk everything
#   - MAP: tiny LLM calls per chunk (return NONE unless answer is present)
#   - REDUCE: synthesize final answer with citations
# -----------------------------------------------------------------------------

from __future__ import annotations

import io, os, re, math, json, string, logging, warnings, mimetypes
from pathlib import Path
from collections import Counter
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Any, Dict, List, Optional, Tuple

import requests

# External connectors (your project)
from graph_api import get_file_with_download_url
from google_drive_api import ensure_google_access_token
from dropbox_api import ensure_dropbox_access_token, TMP_LINK_URL
from box_api import ensure_box_access_token

log = logging.getLogger(__name__)

# =========================
# Ollama via langchain_ollama
# =========================
try:
    from langchain_ollama import ChatOllama
except Exception:  # pragma: no cover
    ChatOllama = None  # type: ignore

OLLAMA_MODEL       = os.getenv("OLLAMA_MODEL", "qwen2:7b")
OLLAMA_NUM_CTX     = int(os.getenv("OLLAMA_NUM_CTX", "20000"))
OLLAMA_NUM_PREDICT = int(os.getenv("OLLAMA_NUM_PREDICT", "5000"))
OLLAMA_TEMPERATURE = float(os.getenv("OLLAMA_TEMPERATURE", "0"))
LLM_DEBUG          = bool(int(os.getenv("SC_LLM_DEBUG", "0")))

_llm = None
if ChatOllama is not None:
    _llm = ChatOllama(
        model=OLLAMA_MODEL,
        temperature=OLLAMA_TEMPERATURE,
        num_ctx=OLLAMA_NUM_CTX,
        num_predict=OLLAMA_NUM_PREDICT,
    )

def _llm_respond(prompt: str) -> str:
    """Single-shot call to ChatOllama; returns plain text."""
    if _llm is None:
        log.error("ChatOllama is not available. Please install langchain-ollama and run Ollama.")
        return ""
    try:
        out = _llm.invoke(prompt)
        txt = getattr(out, "content", None) or str(out)
        return (txt or "").strip()
    except Exception as e:
        if LLM_DEBUG:
            log.exception("Ollama .invoke failed: %s", e)
        try:
            txt = _llm.predict(prompt)
            return (txt or "").strip()
        except Exception as e2:
            log.exception("Ollama .predict failed: %s", e2)
            return ""

def _answer_general_query(prompt: str) -> str:
    sys = (
        "You are a precise assistant. Be concise. Avoid speculation. "
        "When asked to extract data from provided passages, only use that data."
    )
    return _llm_respond(f"system: {sys}\nuser: {prompt}")

# =========================
# Env knobs (retrieval, OCR, map-reduce)
# =========================
MAX_BYTES = int(os.getenv("SC_MAX_BYTES", "10485760"))          # 10 MB per file
MAX_CHARS_LLM = int(os.getenv("SC_MAX_CHARS_LLM", "15000"))

CANDIDATES_BY_NAME = int(os.getenv("SC_CANDIDATES_BY_NAME", "4"))
CANDIDATES_FALLBACK = int(os.getenv("SC_CANDIDATES_FALLBACK", "4"))
CANDIDATE_POOL_CAP = int(os.getenv("SC_CANDIDATE_POOL_CAP", "12"))

SC_SCAN_ALL = os.getenv("SC_SCAN_ALL", "false").lower() == "true"
SC_MAX_WORKERS = int(os.getenv("SC_MAX_WORKERS", "6"))

W_FILENAME_TOKEN   = float(os.getenv("SC_W_FILENAME_TOKEN", "1.0"))
W_FILENAME_PHRASE  = float(os.getenv("SC_W_FILENAME_PHRASE", "4.0"))
W_COMBINE_FILENAME = float(os.getenv("SC_W_COMBINE_FILENAME", "0.35"))
W_COMBINE_CONTENT  = float(os.getenv("SC_W_COMBINE_CONTENT",  "0.65"))

CHUNK_SIZE    = int(os.getenv("SC_CHUNK_SIZE", "3000"))
CHUNK_OVERLAP = int(os.getenv("SC_CHUNK_OVERLAP", "300"))
TOP_CHUNKS    = int(os.getenv("SC_TOP_CHUNKS", "4"))

ARCHIVE_EXTS = {e.strip().lower() for e in os.getenv("SC_ARCHIVES", "zip,rar,7z").split(",") if e.strip()}

ENABLE_OCR           = os.getenv("SC_ENABLE_OCR", "true").lower() == "true"
OCR_LANGS            = os.getenv("SC_OCR_LANGS", "eng")
OCR_MAX_IMG_PIXELS   = int(os.getenv("SC_OCR_MAX_IMG_PIXELS", "20000000"))
PDF_OCR_MAX_PAGES    = int(os.getenv("SC_PDF_OCR_MAX_PAGES", "6"))
PDF_OCR_DPI          = int(os.getenv("SC_PDF_OCR_DPI", "200"))

SC_LOG_CONTEXT = os.getenv("SC_LOG_CONTEXT", "false").lower() == "true"

# Exhaustive mode
SC_ALL_TO_LLM      = os.getenv("SC_ALL_TO_LLM", "false").lower() == "true"
SC_ALL_MAX_FILES   = int(os.getenv("SC_ALL_MAX_FILES", "200"))
SC_ALL_MAX_CHARS   = int(os.getenv("SC_ALL_MAX_CHARS", "800000"))
SC_ALL_FILE_CHARS  = int(os.getenv("SC_ALL_FILE_CHARS", "120000"))
SC_MAP_CHUNK_SIZE  = int(os.getenv("SC_MAP_CHUNK_SIZE", "1800"))
SC_MAP_CHUNK_OVER  = int(os.getenv("SC_MAP_CHUNK_OVER", "150"))
SC_MAP_TOPK_JSON   = int(os.getenv("SC_MAP_TOPK_JSON", "3"))
SC_MAP_WORKERS     = int(os.getenv("SC_MAP_WORKERS", "6"))

# =========================
# Stopwords & query keywords
# =========================
_EXTRA_DIALOG = """
i you your yours me my mine we our ours us please kindly could would should can
tell show find get give list provide share help need want looking
"""
_BASE_STOP = set("""
a an the and or to for of in on at by with from as is are was were be been being
this that these those into over under out up down then else what which who whom
whose where when why how about regarding related contents content file document
pdf ppt pptx doc docx xls xlsx csv link open show give find get see read phone number
""".split()) | set(_EXTRA_DIALOG.split())
_EXTRA_STOP = {w.strip() for w in os.getenv("SC_EXTRA_STOPWORDS", "").split(",") if w.strip()}
STOPWORDS = _BASE_STOP | _EXTRA_STOP

def _normalize_space(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").strip())

def _keywords(q: str, limit: int = 12) -> List[str]:
    q = (q or "").lower()
    phrases = [m.group(1).strip() for m in re.finditer(r'"([^"]+)"', q)]
    residual = re.sub(r'"[^"]+"', " ", q)
    residual = residual.translate(str.maketrans({c: " " for c in string.punctuation if c not in {"-","_","."}}))
    tokens = [t for t in residual.split() if t and (len(t) > 1 or t.isdigit()) and t not in STOPWORDS]
    # dedupe preserve order
    seen, out = set(), []
    for t in (phrases + tokens):
        if t not in seen:
            seen.add(t); out.append(t)
    return out[:limit]

# =========================
# Upload root resolution (for local files)
# =========================
def _get_upload_root() -> str:
    """
    Try Flask config first, then env, then ./uploads as a last resort.
    """
    root = None
    try:
        from flask import current_app
        if current_app:
            root = current_app.config.get("UPLOAD_ROOT")
    except Exception:
        root = None
    return root or os.getenv("UPLOAD_ROOT") or os.path.join(os.getcwd(), "uploads")

def _resolve_local_path(f: Dict[str, Any]) -> Optional[str]:
    """
    Resolve a local file path from common fields:
      - abs_path / absolute_path / local_path / full_path
      - relative_path / path (joined under UPLOAD_ROOT)
      - url or webUrl that begins with /uploads/...
    """
    # Absolute style keys
    for k in ("abs_path", "absolute_path", "local_path", "full_path"):
        p = (f.get(k) or "").strip()
        if p and os.path.isabs(p) and os.path.exists(p):
            return p

    # Relative style keys
    up_root = _get_upload_root()
    for k in ("relative_path", "path"):
        rel = (f.get(k) or "").strip()
        if rel:
            # If absolute was dropped here, keep it
            if os.path.isabs(rel) and os.path.exists(rel):
                return rel
            # Join under uploads
            cand = os.path.join(up_root, rel.lstrip(r"/\ "))
            if os.path.exists(cand):
                return cand

    # URL style: /uploads/...
    for k in ("url", "webUrl", "web_url"):
        u = (f.get(k) or "").strip()
        if u.startswith("/uploads/"):
            cand = os.path.join(up_root, u[len("/uploads/"):].lstrip(r"/\ "))
            if os.path.exists(cand):
                return cand

    return None

def _dl_local(f: Dict[str, Any]) -> Tuple[bytes, Optional[str]]:
    p = _resolve_local_path(f)
    if not p:
        return b"", None
    try:
        # Clamp read size
        with open(p, "rb") as fh:
            data = fh.read(MAX_BYTES)
        mime = mimetypes.guess_type(p)[0]
        log.info("ðŸ“„ reading local upload: %s (%s)", os.path.basename(p), mime or "application/octet-stream")
        return data, mime
    except Exception as e:
        log.exception("Failed to read local upload: %s", e)
        return b"", None

# =========================
# Provider downloads
# =========================
def _dl_ms(f: dict, ms_token: Optional[str]) -> Tuple[bytes, Optional[str]]:
    if not ms_token:
        return b"", None
    meta = get_file_with_download_url(f["parentReference"]["driveId"], f["id"], ms_token)
    if not meta or "@microsoft.graph.downloadUrl" not in meta:
        return b"", None
    url = meta["@microsoft.graph.downloadUrl"]
    r = requests.get(url, stream=True, timeout=30)
    return r.raw.read(MAX_BYTES, decode_content=True), (meta.get("file", {}) or {}).get("mimeType") or r.headers.get("Content-Type")

def _dl_gdrive(uid: int, file_id: str) -> Tuple[bytes, Optional[str]]:
    at = ensure_google_access_token(uid)
    if not at:
        return b"", None
    r = requests.get(f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media",
                     headers={"Authorization": f"Bearer {at}"}, stream=True, timeout=30)
    if not r.ok:
        return b"", None
    return r.raw.read(MAX_BYTES, decode_content=True), r.headers.get("Content-Type")

def _dl_dropbox(uid: int, path_display: str) -> Tuple[bytes, Optional[str]]:
    at = ensure_dropbox_access_token(uid)
    if not at or not path_display:
        return b"", None
    r = requests.post(
        TMP_LINK_URL,
        headers={"Authorization": f"Bearer {at}", "Content-Type": "application/json"},
        json={"path": path_display},
        timeout=20,
    )
    if not r.ok:
        return b"", None
    link = (r.json() or {}).get("link")
    if not link:
        return b"", None
    r2 = requests.get(link, stream=True, timeout=30)
    return r2.raw.read(MAX_BYTES, decode_content=True), r2.headers.get("Content-Type")

def _dl_box(uid: int, file_id: str) -> Tuple[bytes, Optional[str]]:
    at = ensure_box_access_token(uid)
    if not at:
        return b"", None
    r = requests.get(f"https://api.box.com/2.0/files/{file_id}/content",
                     headers={"Authorization": f"Bearer {at}"}, stream=True, timeout=30, allow_redirects=True)
    if not r.ok:
        return b"", None
    return r.raw.read(MAX_BYTES, decode_content=True), r.headers.get("Content-Type")

def download_file_bytes(uid: int, ms_token: Optional[str], f: Dict[str, Any]) -> Tuple[bytes, Optional[str], str]:
    """
    Unified byte loader for:
      - MS Graph / Google Drive / Dropbox / Box
      - Local uploads (source: 'upload'/'local'/'file' or resolvable local path/url)
    """
    name = (f.get("name") or f.get("title") or "file")
    src = (f.get("source") or "").lower().strip()

    # Local upload short-circuit: either explicit source OR we can resolve a path
    local_possible = bool(_resolve_local_path(f))
    if src in ("upload", "local", "file", "localfile", "upload_local") or local_possible:
        data, mime = _dl_local(f)
        return (data or b""), mime, name

    try:
        if src in ("", "ms_graph", "microsoft", "sharepoint", "onedrive"):
            data, mime = _dl_ms(f, ms_token)
        elif src == "google_drive":
            data, mime = _dl_gdrive(uid, f.get("id"))
        elif src == "dropbox":
            data, mime = _dl_dropbox(uid, f.get("path_display") or f.get("path_lower") or "")
        elif src == "box":
            data, mime = _dl_box(uid, f.get("id"))
        else:
            # Fallback try-local if unknown source but path resolves
            if local_possible:
                data, mime = _dl_local(f)
            else:
                return b"", None, name
        return (data or b""), mime, name
    except Exception:
        return b"", None, name

# =========================
# OCR helpers
# =========================
def _try_import_pytesseract():
    try:
        import pytesseract
        from PIL import Image, ImageOps
        return pytesseract, Image, ImageOps
    except Exception:
        return None, None, None

def _ocr_image_bytes(data: bytes) -> str:
    if not ENABLE_OCR:
        return ""
    pytesseract, Image, ImageOps = _try_import_pytesseract()
    if not pytesseract or not Image:
        return ""
    try:
        Image.MAX_IMAGE_PIXELS = OCR_MAX_IMAGE_PIXELS
    except NameError:
        # Back-compat: earlier var name
        from PIL import Image as _Image
        _Image.MAX_IMAGE_PIXELS = OCR_MAX_IMG_PIXELS
    try:
        im = Image.open(io.BytesIO(data))
        if im.mode not in ("L", "RGB", "RGBA"):
            im = im.convert("RGB")
        im = ImageOps.autocontrast(im.convert("L"))
        return (pytesseract.image_to_string(im, lang=OCR_LANGS, config="--oem 3 --psm 6") or "").strip()
    except Exception:
        return ""

def _pdf_ocr(data: bytes) -> str:
    if not ENABLE_OCR:
        return ""
    try:
        import fitz  # PyMuPDF
    except Exception:
        return ""
    pytesseract, Image, ImageOps = _try_import_pytesseract()
    if not pytesseract or not Image:
        return ""
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        out, pages = [], min(len(doc), PDF_OCR_MAX_PAGES)
        scale = PDF_OCR_DPI / 72.0
        mat = fitz.Matrix(scale, scale)
        for i in range(pages):
            pix = doc[i].get_pixmap(matrix=mat)
            o = _ocr_image_bytes(pix.tobytes("png"))
            if o:
                out.append(o)
        return "\n\n".join(out)
    except Exception:
        return ""

# =========================
# Text extraction per type
# =========================
def _pdf(data: bytes) -> str:
    text = ""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        text = "\n".join((p.extract_text() or "") for p in reader.pages[:200])
    except Exception:
        text = ""
    if not text or len(text.strip()) < 50:
        ocr = _pdf_ocr(data)
        if ocr:
            return ocr
    return text

def _docx(data: bytes) -> str:
    try:
        import docx
        d = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in d.paragraphs)
    except Exception:
        return ""

def _pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        txt = []
        for slide in prs.slides[:200]:
            for sh in slide.shapes:
                if hasattr(sh, "text"):
                    txt.append(sh.text)
        return "\n".join(txt)
    except Exception:
        return ""

def _xlsx(data: bytes) -> str:
    """
    Robust, table-aware extraction for Excel/CSV:
      - Pandas header=0 first; if empty, retry header=None
      - Preserve strings (IDs/phones), emit '|' separated text per sheet
      - Fallback to openpyxl keeping structure
    """
    try:
        import pandas as pd

        def _to_text(df_dict: "dict[str,pd.DataFrame]") -> str:
            parts = []
            for sname, df in list(df_dict.items())[:8]:
                if len(df) > 2000:
                    df = df.iloc[:2000, :]
                parts.append(f"# Sheet: {sname}\n" + df.to_csv(sep="|", index=False))
            return "\n\n".join(parts)

        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
            h0 = pd.read_excel(io.BytesIO(data), sheet_name=None, engine="openpyxl", dtype=str, na_filter=False)

        def _has_data(df): return not df.dropna(how="all").empty
        if any(_has_data(df) for df in h0.values()):
            return _to_text(h0)

        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
            hN = pd.read_excel(io.BytesIO(data), sheet_name=None, engine="openpyxl", dtype=str, na_filter=False, header=None)
        hN = {s: df.rename(columns=lambda i: f"col_{i}") for s, df in hN.items()}
        return _to_text(hN)

    except Exception:
        pass

    # openpyxl fallback
    try:
        from openpyxl import load_workbook
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets[:8]:
            rows = []
            for ridx, r in enumerate(ws.iter_rows(values_only=True)):
                vals = ["" if v is None else str(v) for v in list(r)[:60]]
                rows.append(" | ".join(vals))
                if ridx >= 2000: break
            parts.append(f"# Sheet: {ws.title}\n" + "\n".join(rows))
        return "\n\n".join(parts)
    except Exception:
        return ""

def bytes_to_text(data: bytes, name: str, mime: Optional[str]) -> str:
    ext = (os.path.splitext(name)[1] or "").lower()
    m = (mime or "").lower()

    if m.startswith("text/") or ext in {".txt", ".md", ".csv"}:
        for enc in ("utf-8", "latin-1"):
            try:
                return data.decode(enc, errors="ignore")
            except Exception:
                continue
        return ""

    if ext == ".pdf" or m.endswith("/pdf"): return _pdf(data)
    if ext == ".docx": return _docx(data)
    if ext == ".pptx": return _pptx(data)
    if ext in {".xlsx", ".xlsm"}: return _xlsx(data)

    if m.startswith("image/") or ext in {".png",".jpg",".jpeg",".webp",".bmp",".tiff",".tif"}:
        return _ocr_image_bytes(data)

    return ""

# =========================
# Candidate selection & scoring (BM25)
# =========================
def _is_archive(name: str) -> bool:
    return os.path.splitext(name or "")[1].lower().lstrip(".") in ARCHIVE_EXTS

def _filename_score(name: str, q: str) -> float:
    name = (name or "").lower()
    s = 0.0
    for m in re.finditer(r'"([^"]+)"', q.lower()):
        phr = m.group(1).strip()
        if phr and phr in name:
            s += W_FILENAME_PHRASE
    for k in _keywords(q):
        if k in name:
            s += W_FILENAME_TOKEN
    return s

def _build_candidate_pool(hits: List[Dict[str, Any]], q: str) -> List[Dict[str, Any]]:
    if not hits:
        return []
    if SC_SCAN_ALL:
        return list(hits)
    scored = [(_filename_score(h.get("name") or "", q), i, h) for i, h in enumerate(hits)]
    scored.sort(key=lambda t: (-t[0], t[1]))
    pool = [h for _, _, h in scored[:CANDIDATES_BY_NAME]]
    for h in hits[:CANDIDATES_FALLBACK]:
        if h not in pool: pool.append(h)
    return pool[:CANDIDATE_POOL_CAP]

def _tokenize_for_match(text: str) -> List[str]:
    t = re.sub(r"[^0-9a-zA-Z]+", " ", text.lower())
    return t.split()

def _bm25_prepare(docs: List[str], kws: List[str]):
    tok_docs = [_tokenize_for_match(d) for d in docs]
    lows     = [d.lower() for d in docs]
    dls = [len(td) for td in tok_docs]
    avgdl = max(1.0, sum(dls) / max(1, len(dls)))
    counters = [Counter(td) for td in tok_docs]

    N = len(tok_docs)
    df = {}
    for k in kws:
        k_l = k.lower()
        df[k] = sum(1 for td, low in zip(tok_docs, lows)
                    if (k_l in td) or (" " in k_l and k_l in low))
    idf = {k: math.log((N - df[k] + 0.5) / (df[k] + 0.5) + 1.0) for k in kws}
    return counters, dls, avgdl, idf, lows, df, N

def _bm25_score(counter: Counter, low_text: str, dl: int, avgdl: float,
                kws: List[str], idf: Dict[str, float],
                k1: float = 1.5, b: float = 0.75) -> float:
    s = 0.0
    for k in kws:
        k_l = k.lower()
        tf = low_text.count(k_l) if " " in k_l else counter.get(k_l, 0)
        if tf <= 0: continue
        boost = 2.0 if " " in k else 1.0
        denom = tf + k1 * (1 - b + b * (dl / avgdl))
        s += idf.get(k, 0.0) * ((tf * (k1 + 1)) / max(1e-9, denom)) * boost
    return s

def _fetch_text_for_file(uid: int, ms_token: Optional[str], f: dict, user_query: str):
    name = f.get("name") or ""
    log.info("ðŸ”Ž scanning candidate: %s", name)
    if _is_archive(name): return None
    data, mime, _ = download_file_bytes(uid, ms_token, f)
    if not data: return None
    txt = bytes_to_text(data, name, mime)
    if not txt: return None
    return (f, txt)

def _pick_best_by_name_and_content(uid: int,
                                   ms_token: Optional[str],
                                   hits: List[Dict[str, Any]],
                                   user_query: str) -> Tuple[Optional[Dict[str, Any]], Optional[str]]:
    if not hits: return None, None
    pool = _build_candidate_pool(hits, user_query)
    log.info("ðŸ“š candidate pool size: %d (SC_SCAN_ALL=%s)", len(pool), SC_SCAN_ALL)

    docs: List[Tuple[Dict[str, Any], str]] = []
    with ThreadPoolExecutor(max_workers=SC_MAX_WORKERS) as ex:
        futures = [ex.submit(_fetch_text_for_file, uid, ms_token, f, user_query) for f in pool]
        for fut in as_completed(futures):
            res = fut.result()
            if res: docs.append(res)

    if not docs:
        for f in pool:
            if not _is_archive(f.get("name") or ""): return f, None
        return pool[0], None

    kws = _keywords(user_query)
    texts = [t for _, t in docs]
    counters, dls, avgdl, idf, lows, df, N = _bm25_prepare(texts, kws)

    cut = max(1, math.ceil(N / 3))
    rare_terms = {k for k in kws if 0 < df.get(k, 0) <= cut}

    best, best_txt, best_score = None, None, float("-inf")
    for (f, txt), c, dl, low in zip(docs, counters, dls, lows):
        if rare_terms and not any((((" " in k) and (k.lower() in low)) or (c.get(k.lower(), 0) > 0)) for k in rare_terms):
            continue
        s_name = _filename_score(f.get("name") or "", user_query)
        s_text = _bm25_score(c, low, dl, avgdl, kws, idf)
        score = W_COMBINE_FILENAME * s_name + W_COMBINE_CONTENT * s_text
        if score > best_score:
            best, best_txt, best_score = f, txt, score

    if best is None:
        for (f, txt), c, dl, low in zip(docs, counters, dls, lows):
            s_name = _filename_score(f.get("name") or "", user_query)
            s_text = _bm25_score(c, low, dl, avgdl, kws, idf)
            score = W_COMBINE_FILENAME * s_name + W_COMBINE_CONTENT * s_text
            if score > best_score:
                best, best_txt, best_score = f, txt, score

    return best, best_txt

# =========================
# Passage narrowing & LLM prompts
# =========================
def _split_chunks(txt: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[Tuple[int, str]]:
    txt = _normalize_space(txt)
    if not txt: return []
    chunks, i, n = [], 0, len(txt)
    while i < n:
        j = min(n, i + size)
        chunks.append((i, txt[i:j]))
        if j == n: break
        i = max(0, j - overlap)
    return chunks

def _score_chunk(ch: str, q: str) -> float:
    low = ch.lower()
    s = 0.0
    for k in _keywords(q):
        if not k: continue
        s += low.count(k.lower())
    return s / max(300.0, float(len(ch)))

def select_relevant_passages(full_text: str, user_query: str) -> List[Tuple[int, str]]:
    if not full_text: return []
    chunks = _split_chunks(full_text)
    scored = [(off, ch, _score_chunk(ch, user_query)) for (off, ch) in chunks]
    scored.sort(key=lambda t: t[2], reverse=True)
    top = [(off, ch) for (off, ch, _) in scored[:TOP_CHUNKS]]
    return top or (chunks[:1] if chunks else [])

def _make_llm_prompt(user_query: str, filename: str, passages: List[Tuple[int, str]]) -> str:
    merged, total = [], 0
    for _, ch in passages:
        if total + len(ch) > MAX_CHARS_LLM: break
        merged.append(ch); total += len(ch)
    context = "\n\n---\n\n".join(merged)
    if SC_LOG_CONTEXT: log.info("LLM context chars: %d", len(context))
    return (
        "You are a precise document analyst.\n"
        "Using ONLY the provided passages from the document, do ALL of the following:\n"
        "1) Extract exactly the information the user asked for. If a specific value exists, return it verbatim.\n"
        "2) Provide a short, clear answer in Markdown, focused only on the requested content.\n"
        "3) If the answer cannot be confirmed from the passages, say so and suggest the closest relevant details.\n\n"
        f"User request:\n{user_query}\n\n"
        f"Document: {filename}\n\n"
        "Passages:\n"
        f"{context}"
    )

def format_summary_with_llm(user_query: str, filename: str, passages: List[Tuple[int, str]]) -> str:
    prompt = _make_llm_prompt(user_query, filename, passages)
    return _answer_general_query(prompt)

# =========================
# Exhaustive Mapâ†’Reduce (scan ALL files)
# =========================
def _cap_text(s: str, n: int) -> str:
    s = s or ""
    return s if len(s) <= n else s[:n]

def _file_to_snippets(file_dict: dict, text: str) -> List[Dict[str, Any]]:
    chunks = _split_chunks(text, size=SC_MAP_CHUNK_SIZE, overlap=SC_MAP_CHUNK_OVER)
    out = []
    fid = file_dict.get("id")
    name = file_dict.get("name") or "file"
    for off, ch in chunks:
        out.append({"file": file_dict, "file_id": fid, "name": name, "offset": off, "text": ch})
    return out

_MAP_PROMPT = (
    "You are an accurate information extractor. Use ONLY the passage to answer.\n"
    "Task: If the passage contains a definite answer to the user's request, return it; otherwise return NONE.\n"
    "User request:\n{question}\n\n"
    "Passage:\n{passage}\n\n"
    "Return a JSON object with keys exactly: found (true/false), answer (string), reason (short).\n"
)

def _ask_map_llm(question: str, passage: str) -> Dict[str, Any]:
    # Ask for JSON; tolerate models that add extra text
    raw = _llm_respond(_MAP_PROMPT.format(question=question, passage=passage))
    m = re.search(r"\{.*\}", raw, flags=re.S)
    try:
        obj = json.loads(m.group(0)) if m else {}
    except Exception:
        obj = {}
    found = bool(obj.get("found") is True or str(obj.get("found")).lower() == "true")
    ans = (obj.get("answer") or "").strip()
    return {"found": found, "answer": ans, "raw": raw}

def _score_snippet(ans: str, q: str) -> float:
    if not ans: return 0.0
    low = ans.lower(); score = 0.0
    for k in _keywords(q): score += low.count(k.lower())
    return score / max(1.0, len(ans) / 200.0)

def summarize_from_all_hits(uid: int,
                            ms_token: Optional[str],
                            hits: List[Dict[str, Any]],
                            user_query: str) -> Dict[str, Any]:
    if not hits:
        return {"answer": "I couldn't find any files to inspect.", "file": None, "snippets": []}

    total_chars = 0
    kept: List[Tuple[dict, str]] = []

    def _dl(f):
        if _is_archive(f.get("name") or ""): return None
        data, mime, name = download_file_bytes(uid, ms_token, f)
        if not data: return None
        txt = bytes_to_text(data, name, mime)
        if not txt: return None
        return (f, _cap_text(txt, SC_ALL_FILE_CHARS))

    with ThreadPoolExecutor(max_workers=SC_MAX_WORKERS) as ex:
        futures = [ex.submit(_dl, f) for f in hits[:SC_ALL_MAX_FILES]]
        for fut in as_completed(futures):
            res = fut.result()
            if not res: continue
            f, txt = res
            if total_chars + len(txt) > SC_ALL_MAX_CHARS: continue
            kept.append((f, txt)); total_chars += len(txt)

    if not kept:
        return {"answer": "I could not read contents from the files found.", "file": None, "snippets": []}

    passages: List[Dict[str, Any]] = []
    for f, txt in kept:
        passages.extend(_file_to_snippets(f, txt))

    def _map_one(p):
        r = _ask_map_llm(user_query, p["text"])
        r.update({"file": p["file"], "offset": p["offset"], "name": p["name"], "text": p["text"]})
        return r

    hits_map: List[Dict[str, Any]] = []
    with ThreadPoolExecutor(max_workers=SC_MAP_WORKERS) as ex:
        for fut in as_completed([ex.submit(_map_one, p) for p in passages]):
            try:
                res = fut.result()
                if res.get("found") and res.get("answer"):
                    hits_map.append(res)
            except Exception:
                pass

    if not hits_map:
        return summarize_from_hits(uid, ms_token, hits, user_query)

    per_file: Dict[str, List[Dict[str, Any]]] = {}
    for h in hits_map:
        fid = (h["file"] or {}).get("id") or "unknown"
        per_file.setdefault(fid, []).append(h)

    pruned: List[Dict[str, Any]] = []
    for fid, arr in per_file.items():
        arr.sort(key=lambda x: _score_snippet(x["answer"], user_query), reverse=True)
        pruned.extend(arr[:SC_MAP_TOPK_JSON])

    lines = []
    for h in pruned:
        lines.append(f"- Source: {h['name']} (offset {h['offset']})\n  Answer: {h['answer']}")

    reduce_prompt = (
        "You are a precise synthesizer. You will receive candidate answers extracted from multiple documents. "
        "Merge them into ONE short, definitive answer in Markdown. If answers conflict, pick the most specific / most consistent. "
        "If uncertain, say so.\n\n"
        f"Question: {user_query}\n\n"
        "Candidate answers with provenance:\n" + "\n".join(lines)
    )
    final = _answer_general_query(reduce_prompt)

    pruned.sort(key=lambda x: _score_snippet(x["answer"], user_query), reverse=True)
    top_snips = [{"offset": h["offset"], "text": h["text"], "file": h["file"]} for h in pruned[:3]]

    return {"answer": final, "file": None, "snippets": top_snips}

# =========================
# Primary entry point (default mode)
# =========================
def summarize_from_hits(uid: int,
                        ms_token: Optional[str],
                        hits: List[Dict[str, Any]],
                        user_query: str) -> Dict[str, Any]:
    """
    Returns:
      {
        "answer": "<markdown>",
        "file": { ...selected file... } | None,
        "snippets": [{"offset": int, "text": "..."}]
      }
    """
    if SC_ALL_TO_LLM:
        return summarize_from_all_hits(uid, ms_token, hits, user_query)

    f, pre = _pick_best_by_name_and_content(uid, ms_token, hits, user_query)
    if not f:
        return {"answer": "I couldn't find a readable file that matches your request.", "file": None, "snippets": []}

    name = f.get("name") or "file"
    if pre is None:
        data, mime, _ = download_file_bytes(uid, ms_token, f)
        if not data:
            return {"answer": f"I found **{name}**, but couldn't download it from the source.", "file": f, "snippets": []}
        text = bytes_to_text(data, name, mime)
        if not text:
            link = f.get("webUrl") or f.get("web_url") or f.get("url") or ""
            return {"answer": f"I found **{name}**, but couldn't read its contents. {('Open it here: ' + link) if link else ''}", "file": f, "snippets": []}
    else:
        text = pre

    passages = select_relevant_passages(text, user_query)
    answer = format_summary_with_llm(user_query, name, passages)
    snippets = [{"offset": off, "text": snip} for (off, snip) in passages]
    return {"answer": answer, "file": f, "snippets": snippets}
